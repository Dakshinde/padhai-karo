# app.py
import streamlit as st
import time
import io
import json
import PyPDF2
from pptx import Presentation
from docx import Document
from PIL import Image
from datetime import datetime

# PDF generation
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import inch
from reportlab.lib.enums import TA_CENTER
from reportlab.pdfgen import canvas as pdfcanvas

# --- Import AI / extraction utilities from your utils module ---
from utils.gemini_api import (
    generate_quiz_from_topic,
    get_study_resources,
    generate_learning_path,
    extract_topics_from_syllabus,
    generate_quiz_from_context,
    process_syllabus,
    process_pyqs,
    generate_module_question_bank  # returns (parsed_result, raw_text) when debug=True
)

# ---------------------- Helper Functions for Text Extraction ----------------------
def extract_text_from_pdf(file_bytes):
    """Extract text from PDF bytes using PyPDF2."""
    try:
        reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
        text = ""
        for page in reader.pages:
            text += (page.extract_text() or "") + "\n"
        return text
    except Exception as e:
        print(f"extract_text_from_pdf error: {e}")
        return ""

def extract_text_from_pptx(file_bytes):
    """Extract text from PPTX bytes using python-pptx."""
    try:
        presentation = Presentation(io.BytesIO(file_bytes))
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"extract_text_from_pptx error: {e}")
        return ""

def extract_text_from_docx(file_bytes):
    """Extract text from DOCX bytes using python-docx."""
    try:
        document = Document(io.BytesIO(file_bytes))
        text = ""
        for para in document.paragraphs:
            text += para.text + "\n"
        return text
    except Exception as e:
        print(f"extract_text_from_docx error: {e}")
        return ""

# ---------------------- Numbered Canvas for PDF page numbers ----------------------
class NumberedCanvas(pdfcanvas.Canvas):
    """
    Canvas that records page states and writes page numbers & footer on save.
    """
    def __init__(self, *args, footer_text: str = "", **kwargs):
        super().__init__(*args, **kwargs)
        self._saved_page_states = []
        self.footer_text = footer_text

    def showPage(self):
        self._saved_page_states.append(dict(self.__dict__))
        self._startPage()

    def save(self):
        """Add page info to each saved page state, then save."""
        num_pages = len(self._saved_page_states)
        for state in self._saved_page_states:
            self.__dict__.update(state)
            self.draw_page_number(num_pages)
            super().showPage()
        super().save()

    def draw_page_number(self, page_count):
        page_num_text = f"Page {self._pageNumber} of {page_count}"
        self.setFont("Helvetica", 9)
        width, height = A4
        # Page number centered
        self.drawCentredString(width / 2.0, 15, page_num_text)
        # Footer left: branding
        if self.footer_text:
            self.setFont("Helvetica", 9)
            self.drawString(36, 15, self.footer_text)

# ---------------------- PDF Export (continuous layout) ----------------------
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import inch
from reportlab.lib.enums import TA_CENTER
from reportlab.pdfgen import canvas as pdfcanvas
import io
from datetime import datetime

class NumberedCanvas(pdfcanvas.Canvas):
    """Custom canvas to add page numbers and footer text."""
    def __init__(self, *args, footer_text: str = "", **kwargs):
        super().__init__(*args, **kwargs)
        self._saved_page_states = []
        self.footer_text = footer_text

    def showPage(self):
        self._saved_page_states.append(dict(self.__dict__))
        self._startPage()

    def save(self):
        """Add page numbers & footer on every page before saving."""
        num_pages = len(self._saved_page_states)
        for state in self._saved_page_states:
            self.__dict__.update(state)
            self.draw_page_number(num_pages)
            super().showPage()
        super().save()

    def draw_page_number(self, page_count):
        width, height = A4
        self.setFont("Helvetica", 9)
        page_num_text = f"Page {self._pageNumber} of {page_count}"
        self.drawCentredString(width / 2.0, 15, page_num_text)
        if self.footer_text:
            self.drawRightString(width - 36, 15, self.footer_text)

def export_question_bank_pdf(subject_name, question_bank):
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.enums import TA_CENTER
    from reportlab.lib.units import inch
    from reportlab.pdfgen import canvas as pdfcanvas
    import io
    from datetime import datetime

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer, pagesize=A4,
        rightMargin=36, leftMargin=36,
        topMargin=72, bottomMargin=54
    )

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        name="TitleCentered", parent=styles["Title"], alignment=TA_CENTER, spaceAfter=24
    )
    module_style = ParagraphStyle(
        name="ModuleHeading", parent=styles["Heading2"], spaceBefore=12, spaceAfter=12
    )
    body_style = ParagraphStyle(
        name="BodyText", parent=styles["BodyText"], spaceBefore=4, spaceAfter=4
    )
    footer_text = "Generated and created by Padhai Karo"

    story = []

    # Cover page
    story.append(Spacer(1, 1 * inch))
    story.append(Paragraph("Module-wise Question Bank", title_style))
    story.append(Paragraph(f"Subject: <b>{subject_name}</b>", styles["Heading2"]))
    date_str = datetime.now().strftime("%B %d, %Y")
    story.append(Paragraph(f"Generated on: {date_str}", styles["Normal"]))
    story.append(Spacer(1, 0.5 * inch))
    story.append(Paragraph(footer_text, styles["Normal"]))
    story.append(PageBreak())

    # Questions by module
    for module_index, (module_name, questions) in enumerate(question_bank.items(), start=1):
        story.append(Paragraph(f"Module {module_index}: {module_name}", module_style))

        # Sort questions: high importance first, then repetition
        qs_sorted = sorted(
            questions,
            key=lambda q: (0 if q.get("importance", "").lower() == "high" else 1,
                           -int(q.get("repetition_count", 0)))
        )

        for q_index, q in enumerate(qs_sorted, start=1):
            q_text = q.get("question_text", "").strip()
            rep = q.get("repetition_count", 0)
            imp = q.get("importance", "Normal")

            # Serial numbering
            prefix = f"{q_index}. "
            if imp.lower() == "high":
                para_text = f"<b>{prefix}{q_text}</b> — (repeated {rep} times)"
            else:
                para_text = f"• {prefix}{q_text} — (repeated {rep} times)"
            story.append(Paragraph(para_text, body_style))

        story.append(Spacer(1, 0.2 * inch))  # space between modules

    # Custom canvas for footer + page numbers
    class NumberedCanvas(pdfcanvas.Canvas):
        def __init__(self, *args, footer_text: str = "", **kwargs):
            super().__init__(*args, **kwargs)
            self._saved_page_states = []
            self.footer_text = footer_text

        def showPage(self):
            self._saved_page_states.append(dict(self.__dict__))
            self._startPage()

        def save(self):
            num_pages = len(self._saved_page_states)
            for state in self._saved_page_states:
                self.__dict__.update(state)
                self.draw_page_number(num_pages)
                super().showPage()
            super().save()

        def draw_page_number(self, page_count):
            width, _ = A4
            self.setFont("Helvetica", 9)
            self.drawCentredString(width / 2.0, 15, f"Page {self._pageNumber} of {page_count}")
            if self.footer_text:
                self.drawRightString(width - 36, 15, self.footer_text)

    doc.build(story, canvasmaker=lambda *args, **kwargs: NumberedCanvas(*args, footer_text=footer_text, **kwargs))
    buffer.seek(0)
    return buffer

# ---------------------- Streamlit App Setup ----------------------
st.set_page_config(page_title="Padhai Karo", layout="centered")
st.title("Padhai Karo - Your Personalized Engineering Tutor")

# Initialize session state keys if missing
keys_to_init = [
    'quiz_data', 'user_answers', 'quiz_topic', 'start_time',
    'syllabus_topics', 'pyq_question_bank', 'extracted_syllabus_preview'
]
for key in keys_to_init:
    if key not in st.session_state:
        st.session_state[key] = None

# ---------------------- MODULE-WISE PYQ FEATURE ----------------------
st.markdown("---")
st.header("📘 Module-wise PYQ Question Bank Generator")

with st.expander("About this feature", expanded=False):
    st.write(
        "Provide syllabus (upload image/pdf or paste text), optional course objectives, "
        "and 4–6 previous year question papers (PDF or DOCX). The system will extract modules "
        "from the syllabus, parse PYQs, map questions to modules, count repetitions, "
        "and mark high-importance questions."
    )

# Debugging options (toggle)
with st.expander("Options / Debugging", expanded=False):
    show_extracted_preview_by_default = st.checkbox("Automatically show extracted syllabus preview after extraction", value=True)

with st.form("pyq_bank_form"):
    subject_name = st.text_input("Subject Name:", placeholder="e.g., 'Data Structures and Algorithms'")

    syllabus_input_type = st.radio("How would you like to provide the syllabus?",
                                   ["Upload Syllabus File", "Paste Syllabus Text"])

    syllabus_file = None
    pasted_syllabus_text = None

    if syllabus_input_type == "Upload Syllabus File":
        syllabus_file = st.file_uploader("Upload Syllabus (Image or PDF):", type=['png', 'jpg', 'jpeg', 'pdf'])
    else:
        pasted_syllabus_text = st.text_area("Paste your syllabus text here:", height=200)

    # Button to extract & preview (for uploaded syllabus)
    extract_preview_btn = st.form_submit_button("Extract & Preview Syllabus")

    course_objectives = st.text_area("Paste Course Objectives (Optional):", height=150)

    pyq_files = st.file_uploader(
        "Upload Previous Year Question Papers (PDF or DOCX, 4-6 files):",
        type=['pdf', 'docx'],
        accept_multiple_files=True
    )

    generate_btn = st.form_submit_button("Generate Module-wise Question Bank")

# Handle extraction & preview step (user pressed Extract & Preview Syllabus)
if 'extract_preview_btn' in locals() and extract_preview_btn:
    if syllabus_input_type == "Upload Syllabus File":
        if not syllabus_file:
            st.warning("Please upload a syllabus file to extract.")
        else:
            with st.spinner("Extracting text from syllabus..."):
                extracted = process_syllabus(syllabus_file)
                st.session_state.extracted_syllabus_preview = extracted or ""
                if show_extracted_preview_by_default:
                    st.success("Syllabus extracted. Please review and edit if needed.")
    else:
        # pasted text
        if pasted_syllabus_text and pasted_syllabus_text.strip():
            st.session_state.extracted_syllabus_preview = pasted_syllabus_text
            st.success("Syllabus text taken from paste. You can edit it below before generating.")
        else:
            st.warning("Please paste some syllabus text in the box.")

# Show editable preview if available
if st.session_state.extracted_syllabus_preview is not None:
    st.subheader("Editable Syllabus Preview (you can correct OCR mistakes)")
    edited = st.text_area("Edit extracted syllabus text before generating (this will be used by the AI):",
                          value=st.session_state.extracted_syllabus_preview, height=250, key="edited_syllabus_text")
    # Keep session state in sync
    st.session_state.extracted_syllabus_preview = edited

# Handle generation (Generate Module-wise Question Bank button)
if 'generate_btn' in locals() and generate_btn:
    # Basic validations
    if not subject_name:
        st.warning("Please enter a subject name.")
    elif not st.session_state.extracted_syllabus_preview or not st.session_state.extracted_syllabus_preview.strip():
        st.warning("Please extract/paste and confirm syllabus text (use 'Extract & Preview Syllabus').")
    elif not pyq_files or len(pyq_files) < 4 or len(pyq_files) > 6:
        st.warning("Please upload between 4 and 6 PYQ files.")
    else:
        with st.spinner("Processing PYQ files and generating your question bank..."):
            try:
                # Use the edited/preserved syllabus text
                syllabus_text_final = st.session_state.extracted_syllabus_preview

                # Concatenate PYQs text
                pyqs_text = process_pyqs(pyq_files)

                # Generate module-wise question bank via AI (supports debug)
                parsed_result= generate_module_question_bank(syllabus_text_final, pyqs_text, course_objectives)

                if not parsed_result:
                    st.error("Failed to generate the question bank. Try again or check inputs.")
                else:
                    st.session_state.pyq_question_bank = parsed_result

                    st.success(f"Module-wise Question Bank for {subject_name} generated successfully.")
                    # Display modules
                    for module_name, questions in parsed_result.items():
                        questions_sorted = sorted(
                            questions,
                            key=lambda q: (0 if q.get("importance", "").lower() == "high" else 1,
                                           -int(q.get("repetition_count", 0)))
                        )
                        with st.expander(f"📗 {module_name} ({len(questions_sorted)} questions)"):
                            for q in questions_sorted:
                                q_text = q.get("question_text", "").strip()
                                rep = q.get("repetition_count", 0)
                                imp = q.get("importance", "Normal")
                                if imp.lower() == "high":
                                    st.markdown(f"**⭐ {q_text}**  —  (_repeated {rep} times_)")
                                else:
                                    st.markdown(f"- {q_text}  —  (_repeated {rep} times_)")

                    # PDF export (Download button)
                    try:
                        pdf_buffer = export_question_bank_pdf(subject_name, parsed_result)
                        st.download_button(
                            label="📥 Download Question Bank PDF",
                            data=pdf_buffer,
                            file_name=f"{subject_name.replace(' ', '_')}_Question_Bank.pdf",
                            mime="application/pdf"
                        )
                    except Exception as e:
                        st.error(f"Could not create PDF: {e}")


            except Exception as e:
                st.error(f"Error while generating question bank: {e}")

# ---------------------- QUIZ FLOW (Existing App Functionality) ----------------------
# STATE 1: A quiz is active
if st.session_state.quiz_data and st.session_state.user_answers is None:
    st.header(f"Quiz on: {st.session_state.quiz_topic}")
    with st.form("quiz_form"):
        temp_user_answers = {}
        for i, q in enumerate(st.session_state.quiz_data):
            st.subheader(f"Question {i+1}: {q['question_text']}")
            temp_user_answers[i] = st.radio("Choose one:", q['options'], key=f"q{i}", index=None)

        if st.form_submit_button("Submit Answers"):
            st.session_state.user_answers = temp_user_answers
            st.rerun()

# STATE 2: A quiz has been submitted and results are shown
elif st.session_state.user_answers:
    st.header("Your Results")
    quiz_data = st.session_state.quiz_data
    user_answers = st.session_state.user_answers

    score = 0
    incorrect_questions = []
    for i, q in enumerate(quiz_data):
        try:
            if user_answers.get(i, "").strip() == q['correct_answer'].strip():
                score += 1
            else:
                incorrect_questions.append(q)
        except Exception:
            # defensive: if missing keys or trimming issues
            if user_answers.get(i, "") == q.get('correct_answer', ""):
                score += 1
            else:
                incorrect_questions.append(q)

    score_percent = score / len(quiz_data) if len(quiz_data) > 0 else 0
    if score_percent >= 0.9:
        st.success(f"You scored {score} out of {len(quiz_data)}!")
    elif score_percent >= 0.5:
        st.warning(f"You scored {score} out of {len(quiz_data)}!")
    else:
        st.error(f"You scored {score} out of {len(quiz_data)}!")

    st.subheader("Detailed Feedback")
    for i, q in enumerate(quiz_data):
        expanded = (user_answers.get(i, "").strip() != q.get('correct_answer', '').strip())
        with st.expander(f"Question {i+1}: Review", expanded=expanded):
            st.markdown(f"**Question:** {q['question_text']}")
            if user_answers.get(i, "").strip() == q.get('correct_answer', '').strip():
                st.success(f"Your answer: {user_answers.get(i)} (Correct)")
            else:
                st.error(f"Your answer: {user_answers.get(i)} (Incorrect)")
                st.success(f"Correct answer: {q.get('correct_answer', '')}")
            st.info(f"**Explanation:** {q.get('explanation', 'No explanation provided.')}")

    if incorrect_questions:
        incorrect_questions_tuple = tuple(tuple(d.items()) for d in incorrect_questions)
        with st.spinner("Generating your personalized plan..."):
            learning_path = generate_learning_path(st.session_state.quiz_topic, incorrect_questions_tuple)
            study_resources = get_study_resources(st.session_state.quiz_topic, incorrect_questions_tuple)

        if learning_path and "learning_path" in learning_path:
            st.divider()
            st.subheader("Your Personalized Learning Path")
            for i, step in enumerate(learning_path["learning_path"]):
                with st.container():
                    st.markdown(f"**Step {i+1}: {step.get('step_title','')}**")
                    st.markdown(f"**Action:** {step.get('step_details','')}")
                    st.caption(f"**Why:** {step.get('step_rationale','')}")

        if study_resources and "study_plan" in study_resources:
            st.divider()
            st.subheader("Recommended Study Resources")
            for item in study_resources["study_plan"]:
                with st.container():
                    st.markdown(f"**Sub-Topic to Review:** {item.get('sub_topic','')}")
                    st.markdown(f"**How to Study:** {item.get('study_strategy','')}")
                    # include link if available
                    link = item.get("google_search_link")
                    if link:
                        try:
                            st.markdown(f"[Search for this topic]({link})")
                        except Exception:
                            pass

    st.divider()
    if st.button("Start a New Quiz"):
        for key in keys_to_init:
            st.session_state[key] = None
        st.rerun()

# STATE 3: The initial home screen for generating a quiz
else:
    st.header("1. Choose Your Quiz Method")

    topic_tab, syllabus_tab, notes_tab = st.tabs(["By Topic", "From Syllabus", "From My Notes"])

    # --- By Topic ---
    with topic_tab:
        with st.form("topic_quiz_form"):
            st.subheader("Generate a quiz from any topic name")
            topic_input = st.text_input("What topic do you want a quiz on?", placeholder="e.g., 'Red-Black Trees'")
            num_questions_topic = st.slider("Number of Questions:", min_value=5, max_value=20, value=5, key="topic_q")
            quiz_context_topic = st.selectbox("Quiz Purpose:", ("Quick Review / Viva Prep", "Semester Exam Prep"), key="topic_c")
            if st.form_submit_button("Generate Quiz"):
                if topic_input:
                    st.session_state.quiz_topic = topic_input
                    with st.spinner("Generating your quiz..."):
                        st.session_state.quiz_data = generate_quiz_from_topic(topic_input, num_questions_topic, quiz_context_topic)
                        st.rerun()
                else:
                    st.warning("Please enter a topic.")

    # --- From Syllabus ---
    with syllabus_tab:
        st.subheader("Analyze a syllabus to generate quiz topics")
        if st.session_state.syllabus_topics is None:
            syllabus_text_area = st.text_area("Paste your course syllabus here:", height=200)
            if st.button("Analyze Syllabus"):
                if syllabus_text_area:
                    with st.spinner("Analyzing syllabus..."):
                        st.session_state.syllabus_topics = extract_topics_from_syllabus(syllabus_text_area)
                        st.rerun()
                else:
                    st.warning("Please paste your syllabus text.")
        else:
            st.info("Syllabus analyzed! Now configure your quiz below.")
            with st.form("syllabus_quiz_form"):
                selected_topic = st.selectbox("Choose a topic from your syllabus:", options=st.session_state.syllabus_topics)
                num_questions_syllabus = st.slider("Number of Questions:", min_value=5, max_value=20, value=5, key="syllabus_q")
                quiz_context_syllabus = st.selectbox("Quiz Purpose:", ("Quick Review", "Semester Exam Prep"), key="syllabus_c")
                if st.form_submit_button("Generate Quiz from Syllabus Topic"):
                    st.session_state.quiz_topic = selected_topic
                    with st.spinner("Generating your quiz..."):
                        st.session_state.quiz_data = generate_quiz_from_topic(selected_topic, num_questions_syllabus, quiz_context_syllabus)
                        st.rerun()

    # --- From My Notes ---
    with notes_tab:
        st.subheader("Upload your notes to generate a quiz")
        with st.form("notes_quiz_form"):
            uploaded_file = st.file_uploader("Upload your notes (PDF, PPTX, DOCX)", type=['pdf', 'pptx', 'docx'])
            num_questions_notes = st.slider("Number of Questions:", min_value=5, max_value=20, value=5, key="notes_q")
            if st.form_submit_button("Generate Quiz from My Notes"):
                if uploaded_file is not None:
                    extracted_text = ""
                    with st.spinner(f"Reading your file: {uploaded_file.name}..."):
                        file_bytes = uploaded_file.getvalue()
                        try:
                            if uploaded_file.type == "application/pdf":
                                extracted_text = extract_text_from_pdf(file_bytes)
                            elif "presentation" in uploaded_file.type or uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                                extracted_text = extract_text_from_pptx(file_bytes)
                            elif "wordprocessingml" in uploaded_file.type or uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                                extracted_text = extract_text_from_docx(file_bytes)
                        except Exception as e:
                            st.error(f"Could not read the file. Error: {e}")

                    if extracted_text:
                        st.session_state.quiz_topic = f"Document: {uploaded_file.name}"
                        with st.spinner("Generating quiz from document..."):
                            st.session_state.quiz_data = generate_quiz_from_context(extracted_text, num_questions_notes)
                            st.rerun()
                    else:
                        st.warning("Could not extract text from the document.")
                else:
                    st.warning("Please upload a file.")
